"""
Enterprise Risk Quantification & Analytics Engine - Streamlit Dashboard
Main dashboard with multiple tabs for risk analysis and visualization
"""

import io
import sys
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st

# Add src to path
sys.path.insert(0, str(Path(__file__).parent))

from curves import LossExceedanceCurve
from monte_carlo import MonteCarloSimulator
from risk_register import RiskRegister

# Try to import python-pptx for PowerPoint export
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Page configuration
st.set_page_config(
    page_title="Enterprise Risk Analytics Engine",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Custom CSS
st.markdown(
    """
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        color: #2E86AB;
        text-align: center;
        padding: 1rem 0;
    }
    .metric-card {
        background-color: #f0f2f6;
        padding: 1rem;
        border-radius: 0.5rem;
        border-left: 4px solid #2E86AB;
    }
    .stMetric {
        background-color: #ffffff;
        padding: 1rem;
        border-radius: 0.5rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
</style>
""",
    unsafe_allow_html=True,
)


# Initialize session state
if "risk_register" not in st.session_state:
    st.session_state.risk_register = None
if "simulation_results" not in st.session_state:
    st.session_state.simulation_results = None
if "portfolio_stats" not in st.session_state:
    st.session_state.portfolio_stats = None
if "lec_data" not in st.session_state:
    st.session_state.lec_data = None


def load_risk_data(uploaded_file=None, use_sample=False):
    """Load risk register data"""
    rr = RiskRegister()

    try:
        if use_sample:
            # Load sample data
            sample_path = Path(__file__).parent.parent / "data" / "sample_risk_register.csv"
            df = rr.load_from_csv(str(sample_path))
            st.success(f"✅ Loaded {len(df)} risks from sample data")
        elif uploaded_file is not None:
            # Load uploaded file
            if uploaded_file.name.endswith(".csv"):
                df = rr.load_from_csv(uploaded_file)
            elif uploaded_file.name.endswith((".xlsx", ".xls")):
                df = rr.load_from_excel(uploaded_file)
            else:
                st.error("Unsupported file format. Please upload CSV or Excel file.")
                return None
            st.success(f"✅ Loaded {len(df)} risks from {uploaded_file.name}")
        else:
            return None

        st.session_state.risk_register = rr
        return rr
    except Exception as e:
        st.error(f"Error loading risk data: {str(e)}")
        return None


def run_monte_carlo_simulation(risk_register, n_simulations=10000):
    """Run Monte Carlo simulation"""
    with st.spinner(f"Running {n_simulations:,} Monte Carlo simulations..."):
        simulator = MonteCarloSimulator(n_simulations=n_simulations, random_seed=42)

        # Run simulation
        results = simulator.simulate_portfolio(risk_register.get_risks())
        portfolio_stats = simulator.aggregate_portfolio_risk(results)

        # Store in session state
        st.session_state.simulation_results = results
        st.session_state.portfolio_stats = portfolio_stats

        # Calculate Loss Exceedance Curve
        lec = LossExceedanceCurve()
        lec_data = lec.calculate_lec(portfolio_stats["all_simulations"])
        st.session_state.lec_data = lec_data

        st.success(f"✅ Completed {n_simulations:,} simulations successfully!")

    return results, portfolio_stats, lec_data


def display_risk_register_tab():
    """Display Risk Register Overview Tab"""
    st.header("📋 Risk Register Overview")

    if st.session_state.risk_register is None:
        st.info("👆 Please load a risk register from the sidebar to begin analysis.")
        return

    rr = st.session_state.risk_register
    df = rr.get_risks()

    # Summary statistics
    stats = rr.get_summary_statistics()

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Total Risks", stats["total_risks"])
    with col2:
        st.metric("Active Risks", stats["active_risks"])
    with col3:
        st.metric("Avg Likelihood", f"{stats['avg_likelihood']:.2%}")
    with col4:
        st.metric("Total Potential Impact", f"${stats['total_potential_impact']:,.0f}")

    st.markdown("---")

    # Risk breakdown by category
    col1, col2 = st.columns(2)

    with col1:
        st.subheader("Risk Distribution by Category")
        category_counts = df["category"].value_counts()
        fig = px.pie(
            values=category_counts.values,
            names=category_counts.index,
            title="Risks by Category",
            color_discrete_sequence=px.colors.qualitative.Set3,
        )
        st.plotly_chart(fig, use_container_width=True)

    with col2:
        st.subheader("Inherent vs Residual Risk")
        comparison_df = pd.DataFrame(
            {
                "Risk Type": ["Inherent", "Residual"],
                "Average Score": [stats["avg_inherent_score"], stats["avg_residual_score"]],
            }
        )
        fig = px.bar(
            comparison_df,
            x="Risk Type",
            y="Average Score",
            title="Average Risk Scores",
            color="Risk Type",
            color_discrete_map={"Inherent": "#E63946", "Residual": "#2A9D8F"},
        )
        st.plotly_chart(fig, use_container_width=True)

    # Risk matrix
    st.subheader("Risk Heat Map (Likelihood vs Impact)")

    fig = px.scatter(
        df,
        x="impact",
        y="likelihood",
        size="residual_risk_score",
        color="category",
        hover_data=["risk_id", "risk_name", "owner"],
        title="Risk Heat Map",
        labels={"impact": "Impact ($)", "likelihood": "Likelihood (%)"},
        size_max=30,
    )

    # Add risk zones
    fig.add_shape(
        type="rect",
        x0=0,
        y0=0.7,
        x1=df["impact"].max() * 0.3,
        y1=1.0,
        fillcolor="yellow",
        opacity=0.2,
        line_width=0,
    )
    fig.add_shape(
        type="rect",
        x0=df["impact"].max() * 0.7,
        y0=0.7,
        x1=df["impact"].max(),
        y1=1.0,
        fillcolor="red",
        opacity=0.2,
        line_width=0,
    )

    st.plotly_chart(fig, use_container_width=True)

    # Detailed risk table
    st.subheader("Detailed Risk Register")

    # Add filters
    col1, col2, col3 = st.columns(3)
    with col1:
        categories = ["All"] + list(df["category"].unique())
        selected_category = st.selectbox("Filter by Category", categories)
    with col2:
        statuses = ["All"] + list(df["status"].unique())
        selected_status = st.selectbox("Filter by Status", statuses)
    with col3:
        sort_by = st.selectbox(
            "Sort by", ["residual_risk_score", "inherent_risk_score", "impact", "likelihood"]
        )

    # Apply filters
    filtered_df = df.copy()
    if selected_category != "All":
        filtered_df = filtered_df[filtered_df["category"] == selected_category]
    if selected_status != "All":
        filtered_df = filtered_df[filtered_df["status"] == selected_status]

    filtered_df = filtered_df.sort_values(by=sort_by, ascending=False)

    # Display table
    st.dataframe(
        filtered_df[
            [
                "risk_id",
                "risk_name",
                "category",
                "likelihood",
                "impact",
                "inherent_risk_score",
                "residual_risk_score",
                "owner",
                "status",
            ]
        ].style.format(
            {
                "likelihood": "{:.2%}",
                "impact": "${:,.0f}",
                "inherent_risk_score": "${:,.0f}",
                "residual_risk_score": "${:,.0f}",
            }
        ),
        use_container_width=True,
        height=400,
    )


def display_monte_carlo_tab():
    """Display Monte Carlo Simulation Tab"""
    st.header("🎲 Monte Carlo Simulation Results")

    if st.session_state.risk_register is None:
        st.info("👆 Please load a risk register from the sidebar first.")
        return

    # Simulation controls
    col1, col2 = st.columns([2, 1])
    with col1:
        n_simulations = st.slider(
            "Number of Simulations", min_value=1000, max_value=100000, value=10000, step=1000
        )
    with col2:
        if st.button("🚀 Run Simulation", type="primary"):
            run_monte_carlo_simulation(st.session_state.risk_register, n_simulations)

    if st.session_state.simulation_results is None:
        st.info("👆 Click 'Run Simulation' to start the Monte Carlo analysis.")
        return

    results = st.session_state.simulation_results
    portfolio_stats = st.session_state.portfolio_stats

    # Portfolio-level metrics
    st.subheader("📊 Portfolio-Level Risk Metrics")

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Mean Total Loss", f"${portfolio_stats['total_mean_loss']:,.0f}")
    with col2:
        st.metric("Median Total Loss", f"${portfolio_stats['total_median_loss']:,.0f}")
    with col3:
        st.metric("95th Percentile (VaR)", f"${portfolio_stats['total_p95_loss']:,.0f}")
    with col4:
        st.metric("99th Percentile", f"${portfolio_stats['total_p99_loss']:,.0f}")

    st.markdown("---")

    # Distribution plot
    col1, col2 = st.columns(2)

    with col1:
        st.subheader("Portfolio Loss Distribution")
        fig = go.Figure()
        fig.add_trace(
            go.Histogram(
                x=portfolio_stats["all_simulations"],
                nbinsx=50,
                name="Loss Distribution",
                marker_color="#2E86AB",
            )
        )
        fig.add_vline(
            x=portfolio_stats["total_mean_loss"],
            line_dash="dash",
            line_color="red",
            annotation_text="Mean",
        )
        fig.add_vline(
            x=portfolio_stats["total_p95_loss"],
            line_dash="dash",
            line_color="orange",
            annotation_text="95th %ile",
        )
        fig.update_layout(
            xaxis_title="Total Loss ($)", yaxis_title="Frequency", showlegend=False, height=400
        )
        st.plotly_chart(fig, use_container_width=True)

    with col2:
        st.subheader("Risk Contribution by Category")

        # Aggregate by category
        df = st.session_state.risk_register.get_risks()
        results_with_cat = results.merge(df[["risk_id", "category"]], on="risk_id", how="left")

        category_losses = (
            results_with_cat.groupby("category")["mean_loss"].sum().sort_values(ascending=False)
        )

        fig = px.bar(
            x=category_losses.index,
            y=category_losses.values,
            labels={"x": "Category", "y": "Mean Loss ($)"},
            title="Expected Loss by Category",
            color=category_losses.values,
            color_continuous_scale="Reds",
        )
        fig.update_layout(showlegend=False, height=400)
        st.plotly_chart(fig, use_container_width=True)

    # Individual risk results
    st.subheader("Individual Risk Simulation Results")

    display_df = results[
        [
            "risk_id",
            "risk_name",
            "category",
            "mean_loss",
            "median_loss",
            "p95_loss",
            "p99_loss",
            "max_loss",
        ]
    ].sort_values("mean_loss", ascending=False)

    st.dataframe(
        display_df.style.format(
            {
                "mean_loss": "${:,.0f}",
                "median_loss": "${:,.0f}",
                "p95_loss": "${:,.0f}",
                "p99_loss": "${:,.0f}",
                "max_loss": "${:,.0f}",
            }
        ),
        use_container_width=True,
        height=400,
    )

    # Top risks waterfall
    st.subheader("Top 5 Risks - Expected Loss Waterfall")

    top5 = results.nlargest(5, "mean_loss")

    fig = go.Figure(
        go.Waterfall(
            name="Risk Contribution",
            orientation="v",
            measure=["relative"] * 5 + ["total"],
            x=list(top5["risk_name"]) + ["Total"],
            y=list(top5["mean_loss"]) + [top5["mean_loss"].sum()],
            text=[f"${x:,.0f}" for x in list(top5["mean_loss"]) + [top5["mean_loss"].sum()]],
            connector={"line": {"color": "rgb(63, 63, 63)"}},
        )
    )

    fig.update_layout(
        title="Top 5 Risk Contributors to Portfolio Loss", showlegend=False, height=400
    )

    st.plotly_chart(fig, use_container_width=True)


def display_lec_tab():
    """Display Loss Exceedance Curve Tab"""
    st.header("📈 Loss Exceedance Curve (LEC)")

    if st.session_state.lec_data is None:
        st.info("👆 Please run Monte Carlo simulation first to generate Loss Exceedance Curves.")
        return

    lec = LossExceedanceCurve()
    lec.curve_data = st.session_state.lec_data

    # Plot options
    col1, col2 = st.columns([3, 1])
    with col2:
        plot_type = st.radio("Plot Type", ["Plotly (Interactive)", "Matplotlib"])
        show_percentiles = st.checkbox("Show Percentile Markers", value=True)

    with col1:
        if plot_type == "Plotly (Interactive)":
            fig = lec.plot_lec_plotly(add_percentiles=show_percentiles)
            st.plotly_chart(fig, use_container_width=True)
        else:
            fig = lec.plot_lec_matplotlib(add_percentiles=show_percentiles)
            st.pyplot(fig)

    # VaR and CVaR table
    st.subheader("Value at Risk (VaR) and Conditional VaR (CVaR)")

    var_cvar = lec.get_var_cvar(
        st.session_state.portfolio_stats["all_simulations"], confidence_levels=[0.90, 0.95, 0.99]
    )

    col1, col2 = st.columns([2, 1])

    with col1:
        st.dataframe(
            var_cvar.style.format(
                {
                    "confidence_level": "{:.0%}",
                    "confidence_percentage": "{:.0f}%",
                    "var": "${:,.0f}",
                    "cvar": "${:,.0f}",
                }
            ),
            use_container_width=True,
        )

    with col2:
        st.info(
            """
        **VaR (Value at Risk)**: Maximum expected loss at a given confidence level.

        **CVaR (Conditional VaR)**: Average loss beyond the VaR threshold.
        """
        )

    # Return period analysis
    st.subheader("Return Period Analysis")

    col1, col2 = st.columns(2)

    with col1:
        loss_threshold = st.number_input(
            "Enter Loss Threshold ($)",
            min_value=0,
            value=int(st.session_state.portfolio_stats["total_mean_loss"]),
            step=10000,
        )

    with col2:
        # Find exceedance probability for threshold
        lec_data = st.session_state.lec_data
        idx = (lec_data["loss_threshold"] - loss_threshold).abs().idxmin()
        exceedance_prob = lec_data.loc[idx, "exceedance_probability"]
        return_period = lec_data.loc[idx, "return_period"]

        st.metric("Exceedance Probability", f"{exceedance_prob:.2%}")
        st.metric("Return Period (years)", f"{return_period:.1f}")


def display_kpi_dashboard():
    """Display KPI/KRI Dashboard Tab"""
    st.header("📊 KPI/KRI Dashboard")

    if st.session_state.risk_register is None:
        st.info("👆 Please load a risk register from the sidebar first.")
        return

    rr = st.session_state.risk_register
    df = rr.get_risks()

    # Key Risk Indicators
    st.subheader("🎯 Key Risk Indicators (KRIs)")

    col1, col2, col3, col4 = st.columns(4)

    with col1:
        high_risks = len(df[df["residual_risk_score"] > df["residual_risk_score"].quantile(0.75)])
        st.metric("High Priority Risks", high_risks, delta=None)

    with col2:
        avg_likelihood = df["likelihood"].mean()
        st.metric("Avg Risk Likelihood", f"{avg_likelihood:.1%}")

    with col3:
        risk_reduction = (
            (df["inherent_risk_score"].sum() - df["residual_risk_score"].sum())
            / df["inherent_risk_score"].sum()
            * 100
        )
        st.metric("Risk Mitigation %", f"{risk_reduction:.1f}%", delta=f"{risk_reduction:.1f}%")

    with col4:
        if st.session_state.portfolio_stats is not None:
            expected_loss = st.session_state.portfolio_stats["total_mean_loss"]
            st.metric("Expected Annual Loss", f"${expected_loss:,.0f}")
        else:
            st.metric("Expected Annual Loss", "Run simulation")

    st.markdown("---")

    # Inherent vs Residual Risk Comparison
    st.subheader("Inherent vs Residual Risk Analysis")

    col1, col2 = st.columns(2)

    with col1:
        # By risk
        comparison_df = df[["risk_name", "inherent_risk_score", "residual_risk_score"]].copy()
        comparison_df = comparison_df.nlargest(10, "inherent_risk_score")

        fig = go.Figure()
        fig.add_trace(
            go.Bar(
                name="Inherent Risk",
                x=comparison_df["risk_name"],
                y=comparison_df["inherent_risk_score"],
                marker_color="#E63946",
            )
        )
        fig.add_trace(
            go.Bar(
                name="Residual Risk",
                x=comparison_df["risk_name"],
                y=comparison_df["residual_risk_score"],
                marker_color="#2A9D8F",
            )
        )

        fig.update_layout(
            title="Top 10 Risks: Inherent vs Residual",
            xaxis_title="Risk",
            yaxis_title="Risk Score ($)",
            barmode="group",
            height=400,
        )
        st.plotly_chart(fig, use_container_width=True)

    with col2:
        # By category
        category_comparison = (
            df.groupby("category")
            .agg({"inherent_risk_score": "sum", "residual_risk_score": "sum"})
            .reset_index()
        )

        fig = go.Figure()
        fig.add_trace(
            go.Bar(
                name="Inherent Risk",
                x=category_comparison["category"],
                y=category_comparison["inherent_risk_score"],
                marker_color="#E63946",
            )
        )
        fig.add_trace(
            go.Bar(
                name="Residual Risk",
                x=category_comparison["category"],
                y=category_comparison["residual_risk_score"],
                marker_color="#2A9D8F",
            )
        )

        fig.update_layout(
            title="Risk by Category: Inherent vs Residual",
            xaxis_title="Category",
            yaxis_title="Total Risk Score ($)",
            barmode="group",
            height=400,
        )
        st.plotly_chart(fig, use_container_width=True)

    # Risk trends (simulated for demo)
    st.subheader("Risk Trend Analysis")

    # Generate sample trend data
    months = pd.date_range(start="2024-04-01", periods=6, freq="M")
    trend_data = pd.DataFrame(
        {
            "Month": months,
            "Inherent Risk": [
                df["inherent_risk_score"].sum() * (1 + np.random.uniform(-0.1, 0.1))
                for _ in range(6)
            ],
            "Residual Risk": [
                df["residual_risk_score"].sum() * (1 + np.random.uniform(-0.1, 0.1))
                for _ in range(6)
            ],
        }
    )

    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=trend_data["Month"],
            y=trend_data["Inherent Risk"],
            mode="lines+markers",
            name="Inherent Risk",
            line={"color": "#E63946", "width": 3},
        )
    )
    fig.add_trace(
        go.Scatter(
            x=trend_data["Month"],
            y=trend_data["Residual Risk"],
            mode="lines+markers",
            name="Residual Risk",
            line={"color": "#2A9D8F", "width": 3},
        )
    )

    fig.update_layout(
        title="6-Month Risk Trend",
        xaxis_title="Month",
        yaxis_title="Total Risk Score ($)",
        height=400,
    )
    st.plotly_chart(fig, use_container_width=True)

    # Risk appetite vs exposure
    st.subheader("Risk Appetite Analysis")

    col1, col2, col3 = st.columns(3)

    # Simulated risk appetite thresholds
    risk_appetite = df["residual_risk_score"].sum() * 0.8
    risk_tolerance = df["residual_risk_score"].sum() * 1.0
    risk_capacity = df["residual_risk_score"].sum() * 1.2
    current_exposure = df["residual_risk_score"].sum()

    with col1:
        st.metric("Risk Appetite", f"${risk_appetite:,.0f}")
    with col2:
        st.metric(
            "Current Exposure",
            f"${current_exposure:,.0f}",
            delta=f"{((current_exposure - risk_appetite) / risk_appetite * 100):.1f}%",
            delta_color="inverse",
        )
    with col3:
        st.metric("Risk Capacity", f"${risk_capacity:,.0f}")

    # Gauge chart
    fig = go.Figure(
        go.Indicator(
            mode="gauge+number+delta",
            value=current_exposure,
            domain={"x": [0, 1], "y": [0, 1]},
            title={"text": "Risk Exposure vs Appetite"},
            delta={"reference": risk_appetite},
            gauge={
                "axis": {"range": [None, risk_capacity]},
                "bar": {"color": "darkblue"},
                "steps": [
                    {"range": [0, risk_appetite], "color": "lightgreen"},
                    {"range": [risk_appetite, risk_tolerance], "color": "yellow"},
                    {"range": [risk_tolerance, risk_capacity], "color": "red"},
                ],
                "threshold": {
                    "line": {"color": "red", "width": 4},
                    "thickness": 0.75,
                    "value": risk_tolerance,
                },
            },
        )
    )

    fig.update_layout(height=300)
    st.plotly_chart(fig, use_container_width=True)


def generate_powerpoint_deck():
    """Generate PowerPoint presentation with risk analytics"""
    if not HAS_PPTX:
        return None

    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Enterprise Risk Analytics"
        subtitle.text = f"Executive Summary\n{datetime.now().strftime('%B %d, %Y')}"

        # Slide 2: Executive Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Executive Summary"

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        summary_text = generate_executive_summary()
        # Truncate if too long
        tf.text = summary_text[:1000] + "..." if len(summary_text) > 1000 else summary_text

        # Slide 3: KPI/KRI Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Key Risk Indicators"

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame

        if st.session_state.risk_register is not None:
            rr = st.session_state.risk_register
            df = rr.get_risks()

            p = tf.add_paragraph()
            p.text = f"Total Risks: {len(df)}"
            p.level = 0

            p = tf.add_paragraph()
            p.text = f"High Priority Risks: {len(df[df['residual_risk_score'] > df['residual_risk_score'].quantile(0.75)])}"
            p.level = 0

            p = tf.add_paragraph()
            risk_reduction = (
                (df["inherent_risk_score"].sum() - df["residual_risk_score"].sum())
                / df["inherent_risk_score"].sum()
                * 100
            )
            p.text = f"Risk Mitigation: {risk_reduction:.1f}%"
            p.level = 0

            if st.session_state.portfolio_stats is not None:
                p = tf.add_paragraph()
                expected_loss = st.session_state.portfolio_stats["total_mean_loss"]
                p.text = f"Expected Annual Loss: ${expected_loss:,.0f}"
                p.level = 0

        # Save to bytes
        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes.getvalue()

    except Exception as e:
        st.error(f"Error generating PowerPoint: {str(e)}")
        return None


def export_section():
    """Export functionality"""
    st.sidebar.markdown("---")
    st.sidebar.header("📤 Export Options")

    if st.session_state.risk_register is not None:
        # Export risk register
        if st.sidebar.button("📊 Export Risk Register (CSV)"):
            rr = st.session_state.risk_register
            df = rr.get_risks()

            if st.session_state.simulation_results is not None:
                # Include simulation results
                results = st.session_state.simulation_results
                export_df = df.merge(
                    results[["risk_id", "mean_loss", "p95_loss", "p99_loss"]],
                    on="risk_id",
                    how="left",
                )
            else:
                export_df = df

            csv = export_df.to_csv(index=False)
            st.sidebar.download_button(
                label="⬇️ Download CSV",
                data=csv,
                file_name=f"risk_register_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv",
            )

        # Export PowerPoint
        if HAS_PPTX:
            if st.sidebar.button("📊 Generate Executive Deck (PPTX)"):
                with st.spinner("Generating PowerPoint presentation..."):
                    ppt_data = generate_powerpoint_deck()
                    if ppt_data:
                        st.sidebar.download_button(
                            label="⬇️ Download PowerPoint",
                            data=ppt_data,
                            file_name=f"risk_analytics_deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        )
                        st.sidebar.success("✅ PowerPoint generated!")
        else:
            st.sidebar.info(
                "💡 Install python-pptx for PowerPoint export:\npip install python-pptx"
            )

        # Export Text Summary
        if st.sidebar.button("📄 Generate Executive Summary (TXT)"):
            # Create executive summary text
            summary = generate_executive_summary()

            st.sidebar.download_button(
                label="⬇️ Download Summary (TXT)",
                data=summary,
                file_name=f"risk_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                mime="text/plain",
            )


def generate_executive_summary():
    """Generate executive summary report"""
    rr = st.session_state.risk_register
    df = rr.get_risks()
    stats = rr.get_summary_statistics()

    summary = f"""
ENTERPRISE RISK QUANTIFICATION & ANALYTICS
Executive Summary Report
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

{'='*70}
PORTFOLIO OVERVIEW
{'='*70}

Total Risks Identified: {stats['total_risks']}
Active Risks: {stats['active_risks']}
Risk Categories: {stats['categories']}

Average Likelihood: {stats['avg_likelihood']:.2%}
Average Impact: ${stats['avg_impact']:,.0f}

Total Inherent Risk Score: ${stats['avg_inherent_score'] * stats['total_risks']:,.0f}
Total Residual Risk Score: ${stats['avg_residual_score'] * stats['total_risks']:,.0f}
Risk Reduction: {((stats['avg_inherent_score'] - stats['avg_residual_score']) / stats['avg_inherent_score'] * 100):.1f}%

"""

    if st.session_state.portfolio_stats is not None:
        ps = st.session_state.portfolio_stats
        summary += f"""
{'='*70}
MONTE CARLO SIMULATION RESULTS
{'='*70}

Number of Simulations: {ps['n_simulations']:,}

Expected Loss (Mean): ${ps['total_mean_loss']:,.0f}
Median Loss: ${ps['total_median_loss']:,.0f}
Standard Deviation: ${ps['total_std_loss']:,.0f}

Value at Risk (VaR):
  - 90th Percentile: ${ps['total_p90_loss']:,.0f}
  - 95th Percentile: ${ps['total_p95_loss']:,.0f}
  - 99th Percentile: ${ps['total_p99_loss']:,.0f}

Conditional VaR (95%): ${ps['total_cvar_95']:,.0f}
Maximum Simulated Loss: ${ps['total_max_loss']:,.0f}

"""

    # Top risks
    top_risks = df.nlargest(5, "residual_risk_score")
    summary += f"""
{'='*70}
TOP 5 RISKS (by Residual Risk Score)
{'='*70}

"""
    for idx, (_, risk) in enumerate(top_risks.iterrows(), 1):
        summary += f"""
{idx}. {risk['risk_name']} ({risk['risk_id']})
   Category: {risk['category']}
   Likelihood: {risk['likelihood']:.2%}
   Impact: ${risk['impact']:,.0f}
   Residual Risk Score: ${risk['residual_risk_score']:,.0f}
   Owner: {risk['owner']}

"""

    summary += f"""
{'='*70}
RECOMMENDATIONS
{'='*70}

1. High Priority Risks: Focus mitigation efforts on the top 5 identified risks
2. Risk Monitoring: Implement continuous monitoring for high-likelihood events
3. Risk Transfer: Consider insurance or hedging for high-impact, low-likelihood risks
4. Control Effectiveness: Review and enhance risk controls to further reduce residual risk
5. Regular Reviews: Conduct quarterly risk assessments to update the risk register

{'='*70}
END OF REPORT
{'='*70}
"""

    return summary


def main():
    """Main application"""

    # Header
    st.markdown(
        '<div class="main-header">🎯 Enterprise Risk Quantification & Analytics Engine</div>',
        unsafe_allow_html=True,
    )

    # Sidebar
    st.sidebar.title("🔧 Configuration")

    st.sidebar.header("📁 Load Risk Register")

    # Option to use sample data
    use_sample = st.sidebar.checkbox("Use Sample Data", value=True)

    if use_sample:
        if st.sidebar.button("Load Sample Data", type="primary"):
            load_risk_data(use_sample=True)
    else:
        uploaded_file = st.sidebar.file_uploader(
            "Upload Risk Register",
            type=["csv", "xlsx", "xls"],
            help="Upload a CSV or Excel file with risk data",
        )

        if uploaded_file is not None:
            load_risk_data(uploaded_file=uploaded_file)

    # Export section
    export_section()

    # About section
    st.sidebar.markdown("---")
    st.sidebar.header("ℹ️ About")
    st.sidebar.info(
        """
    **Enterprise Risk Analytics Engine**

    A comprehensive risk quantification tool using Monte Carlo simulation and advanced analytics.

    **Features:**
    - Monte Carlo simulation (10k+ iterations)
    - Loss Exceedance Curves
    - Risk heat maps
    - KPI/KRI dashboards
    - Export capabilities
    """
    )

    # Main tabs
    tab1, tab2, tab3, tab4 = st.tabs(
        [
            "📋 Risk Register",
            "🎲 Monte Carlo Simulation",
            "📈 Loss Exceedance Curve",
            "📊 KPI/KRI Dashboard",
        ]
    )

    with tab1:
        display_risk_register_tab()

    with tab2:
        display_monte_carlo_tab()

    with tab3:
        display_lec_tab()

    with tab4:
        display_kpi_dashboard()


if __name__ == "__main__":
    main()
